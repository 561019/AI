from __future__ import annotations

import csv
import json
import mimetypes
import re
import unicodedata
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Protocol

from .models import ParsedContent, ParsedTable, ParsedValue, SourceRef


# Office formats are parsed directly; image/PDF formats are sent to the OCR adapter.
DIRECT_EXTENSIONS = {".csv", ".tsv", ".xlsx", ".xlsm", ".xls", ".docx", ".pptx", ".md", ".markdown", ".json"}
OCR_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".tif", ".tiff", ".bmp", ".pdf"}


def clean_text(value: Any) -> Any:
    """Normalize display/storage text without changing numeric business values."""
    if not isinstance(value, str):
        return value
    value = unicodedata.normalize("NFKC", value).replace("\u00a0", " ").replace("\r\n", "\n").replace("\r", "\n")
    value = "".join(char for char in value if char == "\n" or char >= " ")
    return "\n".join(re.sub(r"[ \t]+", " ", line).strip() for line in value.split("\n")).strip()


def media_type(path: Path) -> str:
    if path.suffix.lower() == ".csv":
        return "text/csv"
    if path.suffix.lower() == ".tsv":
        return "text/tab-separated-values"
    if path.suffix.lower() in {".md", ".markdown"}:
        return "text/markdown"
    if path.suffix.lower() == ".json":
        return "application/json"
    return mimetypes.guess_type(path.name)[0] or "application/octet-stream"


def value_type(value: Any) -> str:
    if isinstance(value, bool):
        return "boolean"
    if isinstance(value, (int, float)):
        return "number"
    if value is None:
        return "null"
    return "text"


class StructuredParser:
    def parse(self, path: Path, file_hash: str) -> ParsedContent:
        suffix = path.suffix.lower()
        if suffix in {".csv", ".tsv"}:
            return self._parse_delimited(path, file_hash, "\t" if suffix == ".tsv" else ",")
        if suffix in {".xlsx", ".xlsm"}:
            return self._parse_excel(path, file_hash)
        if suffix == ".xls":
            return self._parse_xls(path, file_hash)
        if suffix == ".docx":
            return self._parse_docx(path, file_hash)
        if suffix == ".pptx":
            return self._parse_pptx(path, file_hash)
        if suffix in {".md", ".markdown"}:
            return self._parse_markdown(path, file_hash)
        if suffix == ".json":
            return self._parse_json(path, file_hash)
        raise ValueError(f"unsupported direct parser format: {suffix}")

    def _parse_delimited(self, path: Path, file_hash: str, delimiter: str) -> ParsedContent:
        table = ParsedTable(name=path.stem)
        raw_bytes = path.read_bytes()
        decoded: str | None = None
        for encoding in ("utf-8-sig", "gb18030", "utf-16"):
            try:
                decoded = raw_bytes.decode(encoding)
                break
            except UnicodeDecodeError:
                continue
        if decoded is None:
            raise ValueError("CSV/TSV encoding is not supported; use UTF-8, GB18030, or UTF-16")
        for row_idx, row in enumerate(csv.reader(decoded.splitlines(), delimiter=delimiter), start=1):
            for col_idx, raw in enumerate(row, start=1):
                raw = clean_text(raw)
                if raw == "":
                    continue
                table.values.append(ParsedValue(raw_value=raw, value_type="text", source=SourceRef(
                    file_hash, path.name, sheet=path.stem, cell=f"R{row_idx}C{col_idx}", row=row_idx, column=col_idx,
                )))
        return ParsedContent(tables=[table])

    def _parse_excel(self, path: Path, file_hash: str) -> ParsedContent:
        from openpyxl import load_workbook

        workbook = load_workbook(path, data_only=False, read_only=True)
        tables: list[ParsedTable] = []
        try:
            for sheet in workbook.worksheets:
                table = ParsedTable(name=sheet.title)
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value is None:
                            continue
                        raw = clean_text(cell.value)
                        kind = "formula" if isinstance(raw, str) and raw.startswith("=") else value_type(raw)
                        table.values.append(ParsedValue(raw_value=raw, value_type=kind, source=SourceRef(
                            file_hash, path.name, sheet=sheet.title, cell=cell.coordinate, row=cell.row, column=cell.column,
                        )))
                tables.append(table)
        finally:
            workbook.close()
        return ParsedContent(tables=tables)

    def _parse_xls(self, path: Path, file_hash: str) -> ParsedContent:
        import xlrd

        workbook = xlrd.open_workbook(path)
        tables: list[ParsedTable] = []
        for sheet in workbook.sheets():
            table = ParsedTable(name=sheet.name)
            for row_idx in range(sheet.nrows):
                for col_idx in range(sheet.ncols):
                    raw = sheet.cell_value(row_idx, col_idx)
                    if raw == "":
                        continue
                    raw = clean_text(raw)
                    table.values.append(ParsedValue(raw_value=raw, value_type=value_type(raw), source=SourceRef(
                        file_hash, path.name, sheet=sheet.name, cell=f"R{row_idx + 1}C{col_idx + 1}", row=row_idx + 1, column=col_idx + 1,
                    )))
            tables.append(table)
        return ParsedContent(tables=tables)

    def _parse_docx(self, path: Path, file_hash: str) -> ParsedContent:
        from docx import Document

        document = Document(path)
        text_blocks: list[ParsedValue] = []
        tables: list[ParsedTable] = []
        for idx, paragraph in enumerate(document.paragraphs, start=1):
            text = clean_text(paragraph.text)
            if text:
                text_blocks.append(ParsedValue(raw_value=text, source=SourceRef(file_hash, path.name, paragraph=idx)))
        for table_idx, source_table in enumerate(document.tables, start=1):
            table = ParsedTable(name=f"table-{table_idx}")
            for row_idx, row in enumerate(source_table.rows, start=1):
                for col_idx, cell in enumerate(row.cells, start=1):
                    raw = clean_text(cell.text)
                    if raw:
                        table.values.append(ParsedValue(raw_value=raw, source=SourceRef(
                            file_hash, path.name, table=table_idx, row=row_idx, column=col_idx,
                        )))
            tables.append(table)
        return ParsedContent(text_blocks=text_blocks, tables=tables)

    def _parse_pptx(self, path: Path, file_hash: str) -> ParsedContent:
        from pptx import Presentation

        presentation = Presentation(path)
        text_blocks: list[ParsedValue] = []
        tables: list[ParsedTable] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            for shape_index, shape in enumerate(slide.shapes, start=1):
                if getattr(shape, "has_text_frame", False):
                    text = clean_text(shape.text)
                    if text:
                        text_blocks.append(ParsedValue(raw_value=text, source=SourceRef(
                            file_hash, path.name, page=slide_index, paragraph=shape_index,
                        )))
                if getattr(shape, "has_table", False):
                    table = ParsedTable(name=f"slide-{slide_index}-table-{shape_index}")
                    for row_idx, row in enumerate(shape.table.rows, start=1):
                        for col_idx, cell in enumerate(row.cells, start=1):
                            text = clean_text(cell.text)
                            if text:
                                table.values.append(ParsedValue(raw_value=text, source=SourceRef(
                                    file_hash, path.name, page=slide_index, table=shape_index, row=row_idx, column=col_idx,
                                )))
                    tables.append(table)
        return ParsedContent(text_blocks=text_blocks, tables=tables)

    def _parse_markdown(self, path: Path, file_hash: str) -> ParsedContent:
        lines = path.read_text(encoding="utf-8-sig").splitlines()
        text_blocks: list[ParsedValue] = []
        tables: list[ParsedTable] = []
        paragraph: list[str] = []
        table_lines: list[tuple[int, str]] = []

        def flush_paragraph(start: int) -> None:
            nonlocal paragraph
            text = clean_text(" ".join(paragraph))
            if text:
                text_blocks.append(ParsedValue(raw_value=text, source=SourceRef(file_hash, path.name, paragraph=start)))
            paragraph = []

        def flush_table() -> None:
            nonlocal table_lines
            if len(table_lines) < 2:
                for line_no, line in table_lines:
                    text_blocks.append(ParsedValue(raw_value=clean_text(line), source=SourceRef(file_hash, path.name, paragraph=line_no)))
                table_lines = []
                return
            table = ParsedTable(name=f"markdown-table-{len(tables) + 1}")
            for row_idx, (_, line) in enumerate(table_lines, start=1):
                cells = [clean_text(cell) for cell in line.strip().strip("|").split("|")]
                if row_idx == 2 and all(re.fullmatch(r":?-{3,}:?", cell.replace(" ", "")) for cell in cells):
                    continue
                for col_idx, cell in enumerate(cells, start=1):
                    if cell:
                        table.values.append(ParsedValue(raw_value=cell, source=SourceRef(
                            file_hash, path.name, table=len(tables) + 1, row=row_idx, column=col_idx,
                        )))
            if table.values:
                tables.append(table)
            table_lines = []

        paragraph_start = 1
        for line_no, line in enumerate(lines, start=1):
            if "|" in line and line.strip():
                if paragraph:
                    flush_paragraph(paragraph_start)
                table_lines.append((line_no, line))
                continue
            if table_lines:
                flush_table()
            if not line.strip():
                if paragraph:
                    flush_paragraph(paragraph_start)
                continue
            if not paragraph:
                paragraph_start = line_no
            paragraph.append(re.sub(r"^#{1,6}\s*", "", line).strip())
        if table_lines:
            flush_table()
        if paragraph:
            flush_paragraph(paragraph_start)
        return ParsedContent(text_blocks=text_blocks, tables=tables)

    def _parse_json(self, path: Path, file_hash: str) -> ParsedContent:
        payload = json.loads(path.read_text(encoding="utf-8-sig"))
        table = ParsedTable(name=path.stem)

        def walk(value: Any, json_path: str) -> None:
            if isinstance(value, dict):
                for key, child in value.items():
                    walk(child, f"{json_path}.{key}")
            elif isinstance(value, list):
                for index, child in enumerate(value):
                    walk(child, f"{json_path}[{index}]")
            else:
                table.values.append(ParsedValue(raw_value=clean_text(value), value_type=value_type(value), field_name=json_path, source=SourceRef(
                    file_hash, path.name, sheet=path.stem, cell=json_path, row=len(table.values) + 1, column=1,
                )))

        walk(payload, "$")
        return ParsedContent(tables=[table])


class OCRProvider(Protocol):
    def recognize(self, path: Path, file_hash: str) -> ParsedContent: ...


@dataclass
class SidecarOCRProvider:
    """Reads a local ``<original>.ocr.json`` file for offline integration tests."""
    def recognize(self, path: Path, file_hash: str) -> ParsedContent:
        sidecar = Path(str(path) + ".ocr.json")
        if not sidecar.exists():
            raise RuntimeError(f"OCR is not configured and no sidecar file exists: {sidecar}")
        payload = json.loads(sidecar.read_text(encoding="utf-8"))
        blocks: list[ParsedValue] = []
        for item in payload.get("blocks", []):
            bbox = item.get("bbox")
            text = clean_text(item.get("text", ""))
            if text:
                blocks.append(ParsedValue(raw_value=text, confidence=float(item.get("confidence", 0.0)), source=SourceRef(
                    file_hash, path.name, page=int(item.get("page", 1)), bbox=tuple(bbox) if bbox else None,
                )))
        return ParsedContent(text_blocks=blocks)
